"""Generic document extraction for user-uploaded project files.

Unlike the Danone reference loaders (which know each file's exact schema), this
layer extracts plain text + tabular rows from *arbitrary* uploads so the LLM can
ground on whatever the user actually provides. Dispatch is by file extension:

    .pdf            -> pypdf            (page text)
    .pptx           -> python-pptx      (shape text + table rows)
    .docx           -> python-docx      (paragraphs + table rows)
    .xlsx / .xlsm   -> openpyxl         (sheet rows)
    .csv / .txt /.md -> plain read

Returns a uniform ``ExtractResult`` regardless of source type. Extraction never
raises — failures are captured in ``error`` so a bad upload can't break a run.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

# Per-file extraction caps so a huge deck/workbook can't blow the LLM context.
_MAX_CHARS = 20_000
_MAX_TABLE_ROWS = 200
_MAX_SHEET_ROWS = 120


@dataclass
class ExtractResult:
    text: str = ""
    tables: list[list[list[str]]] = field(default_factory=list)
    error: str | None = None

    @property
    def ok(self) -> bool:
        return self.error is None and bool(self.text or self.tables)


def _clip(text: str) -> str:
    return text[:_MAX_CHARS]


def _cell(v: object) -> str:
    return "" if v is None else str(v).strip()


def _row_line(cells: list[str]) -> str:
    """Render one spreadsheet row as a compact ' | '-joined line, dropping the
    trailing empty cells that merged/title layouts leave behind."""
    trimmed = list(cells)
    while trimmed and not trimmed[-1]:
        trimmed.pop()
    return " | ".join(trimmed)


def _extract_pdf(path: Path) -> ExtractResult:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts = [page.extract_text() or "" for page in reader.pages[:60]]
    return ExtractResult(text=_clip("\n".join(parts)))


def _extract_pptx(path: Path) -> ExtractResult:
    from pptx import Presentation

    prs = Presentation(str(path))
    lines: list[str] = []
    tables: list[list[list[str]]] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
            if shape.has_table:
                rows = [[_cell(c.text) for c in row.cells] for row in shape.table.rows]
                tables.append(rows[:_MAX_TABLE_ROWS])
    return ExtractResult(text=_clip("\n".join(lines)), tables=tables)


def _extract_docx(path: Path) -> ExtractResult:
    import docx

    doc = docx.Document(str(path))
    lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    tables: list[list[list[str]]] = []
    for table in doc.tables:
        rows = [[_cell(c.text) for c in row.cells] for row in table.rows]
        tables.append(rows[:_MAX_TABLE_ROWS])
    return ExtractResult(text=_clip("\n".join(lines)), tables=tables)


def _extract_xlsx(path: Path) -> ExtractResult:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    tables: list[list[list[str]]] = []
    text_lines: list[str] = []
    for ws in wb.worksheets:
        rows: list[list[str]] = []
        for r in ws.iter_rows(values_only=True):
            cells = [_cell(c) for c in r]
            if any(cells):
                rows.append(cells)
            if len(rows) >= _MAX_SHEET_ROWS:
                break
        if rows:
            tables.append(rows)
            # Emit every non-empty data row into the text (not just the header
            # row) — grounding callers read only `.text`, so a header-only text
            # would drop the entire workbook body. `_clip` caps total size.
            text_lines.append(f"[sheet: {ws.title}]")
            text_lines.extend(_row_line(row) for row in rows)
    wb.close()
    return ExtractResult(text=_clip("\n".join(text_lines)), tables=tables)


def _extract_plain(path: Path) -> ExtractResult:
    return ExtractResult(text=_clip(path.read_text(encoding="utf-8", errors="replace")))


_DISPATCH = {
    ".pdf": _extract_pdf,
    ".pptx": _extract_pptx,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".xlsm": _extract_xlsx,
    ".csv": _extract_plain,
    ".txt": _extract_plain,
    ".md": _extract_plain,
}

SUPPORTED_EXTENSIONS = tuple(_DISPATCH.keys())

# Audio files aren't text-extractable; they're transcribed by the ASR step
# (task 1.4b) which writes a .txt transcript sidecar that IS extractable.
AUDIO_EXTENSIONS = (".mp3", ".wav", ".m4a", ".aac", ".ogg", ".flac", ".mp4", ".webm")


def is_audio_filename(filename: str) -> bool:
    return Path(filename).suffix.lower() in AUDIO_EXTENSIONS


def extract_document(path: Path) -> ExtractResult:
    """Extract text + tables from a single file. Never raises."""
    if is_audio_filename(path.name):
        return ExtractResult(error="audio file — pending transcription (ASR)")
    fn = _DISPATCH.get(path.suffix.lower())
    if fn is None:
        return ExtractResult(error=f"unsupported file type: {path.suffix or '(none)'}")
    try:
        return fn(path)
    except Exception as exc:  # noqa: BLE001 — extraction must not break a run
        return ExtractResult(error=f"{type(exc).__name__}: {exc}")
